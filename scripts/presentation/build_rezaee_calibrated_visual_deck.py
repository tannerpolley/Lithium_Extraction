from __future__ import annotations

import csv
import json
from pathlib import Path

import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "slides" / "final_rezaee_calibrated_case_study_2026_05_08"
ASSET_DIR = OUT_DIR / "generated_assets"
TEMPLATE_PATH = ROOT / "docs" / "Slides" / "Lithium Extraction with ePC-SAFT.pptx"

TITLE_COLOR = RGBColor(22, 32, 45)
BODY_COLOR = RGBColor(38, 50, 66)
MUTED_COLOR = RGBColor(93, 108, 126)
ACCENT = RGBColor(21, 101, 118)
GREEN = RGBColor(45, 118, 74)
AMBER = RGBColor(169, 104, 45)
PAPER_BG = RGBColor(248, 250, 252)
CARD_STROKE = RGBColor(209, 218, 230)
DEEP_NAVY = RGBColor(15, 27, 42)
TEAL_DARK = RGBColor(13, 83, 98)
TEAL_LIGHT = RGBColor(229, 244, 242)
SLATE_100 = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)


def read_summary() -> dict:
    path = ROOT / "analyses" / "rezaee_2026_pcsaft_epcsaft" / "results" / "uq_surrogate" / "rezaee_tds_li_oa_uq_summary.json"
    return json.loads(path.read_text(encoding="utf-8"))


def read_nominal_row() -> dict[str, str]:
    path = ROOT / "data" / "reference" / "produced_water" / "rezaee_tds_li_oa_uq_predictions.csv"
    with path.open("r", newline="", encoding="utf-8") as handle:
        for row in csv.DictReader(handle):
            if row["case_id"] == "nominal_ms2_clean_li_na":
                return row
    raise RuntimeError("nominal_ms2_clean_li_na not found")


def asset(rel: str) -> Path:
    p = ROOT / rel
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def render_math_png(name: str, formula: str, width: float = 7.2, height: float = 1.1, fontsize: int = 24) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    out = ASSET_DIR / name
    fig = plt.figure(figsize=(width, height), dpi=220)
    fig.patch.set_alpha(0)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.text(0.02, 0.52, formula, fontsize=fontsize, color="#16202D", va="center", ha="left")
    fig.savefig(out, transparent=True, bbox_inches="tight", pad_inches=0.03)
    plt.close(fig)
    return out


def render_pipeline_diagram(summary: dict) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    out = ASSET_DIR / "rezaee_surrogate_pipeline.png"
    fig, ax = plt.subplots(figsize=(9.0, 3.2), dpi=220)
    ax.set_axis_off()
    fig.patch.set_facecolor("#FFFFFF")
    labels = [
        ("Rezaee\nDOE rows", "31 calibrated\nSection 3.2 points"),
        ("Log-D\nsurfaces", r"$\log D_{\mathrm{Li}},\ \log D_{\mathrm{Na}}$"),
        ("Produced-water\nmatrix", "523 TDS-Li-O/A\nscreening rows"),
        ("Process\nhandoff", "Distribution coefficients\nfor staged contacting"),
    ]
    xs = [0.08, 0.33, 0.58, 0.82]
    for i, ((title, sub), x) in enumerate(zip(labels, xs)):
        rect = plt.Rectangle((x - 0.09, 0.32), 0.18, 0.36, facecolor="#F1F5F9", edgecolor="#CBD5E1", linewidth=1.4)
        ax.add_patch(rect)
        ax.text(x, 0.57, title, ha="center", va="center", fontsize=12, fontweight="bold", color="#16202D")
        ax.text(x, 0.41, sub, ha="center", va="center", fontsize=9, color="#526172")
        if i < len(xs) - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.105, 0.50), xytext=(x + 0.105, 0.50),
                        arrowprops=dict(arrowstyle="->", lw=1.6, color="#156576"))
    ax.text(0.08, 0.16, f"Model basis: {summary['model_basis']}", fontsize=9, color="#526172")
    ax.text(0.08, 0.08, "Caveat carried per row: outside original Rezaee Na/Li DOE domain", fontsize=9, color="#A5682D")
    fig.savefig(out, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    return out


def render_process_diagram() -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    out = ASSET_DIR / "process_bridge_diagram.png"
    fig, ax = plt.subplots(figsize=(9.0, 2.9), dpi=220)
    ax.set_axis_off()
    fig.patch.set_facecolor("#FFFFFF")
    steps = [
        ("Smackover\nbrine", "source composition"),
        ("Divalent\npretreatment", r"remove Ca$^{2+}$/Mg$^{2+}$"),
        ("Clean\nLi/Na feed", "surrogate inlet"),
        ("DES/TOPO\nextraction", "distribution\ncoefficients"),
        ("PrOMMiS /\nIDAES", "MSContactor\ncosting"),
    ]
    xs = [0.08, 0.28, 0.48, 0.68, 0.88]
    for i, ((title, sub), x) in enumerate(zip(steps, xs)):
        color = "#E8F5F2" if i in (2, 3) else "#F8FAFC"
        rect = plt.Rectangle((x - 0.078, 0.30), 0.156, 0.38, facecolor=color, edgecolor="#B6C6D7", linewidth=1.2)
        ax.add_patch(rect)
        ax.text(x, 0.57, title, ha="center", va="center", fontsize=11, fontweight="bold", color="#16202D")
        ax.text(x, 0.40, sub, ha="center", va="center", fontsize=8.5, color="#526172")
        if i < len(xs) - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.087, 0.49), xytext=(x + 0.087, 0.49),
                        arrowprops=dict(arrowstyle="->", lw=1.5, color="#156576"))
    ax.text(0.08, 0.12, "Boundary condition: raw multication brine is never treated as the Rezaee extraction feed.", fontsize=9.5, color="#526172")
    fig.savefig(out, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    return out


def render_validation_chart(summary: dict, nominal: dict[str, str]) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    out = ASSET_DIR / "na_li_validation_boundary.png"
    fig, ax = plt.subplots(figsize=(8.2, 3.2), dpi=220)
    fig.patch.set_facecolor("#FFFFFF")
    ax.barh(["Rezaee DOE\ntraining", "MS-2 clean\nprojection"], [summary["training_Na_Li_max"], float(nominal["Na_Li_mass_ratio"])],
            color=["#156576", "#A5682D"], height=0.42)
    ax.set_xlabel("Na/Li mass ratio")
    ax.set_xlim(0, 420)
    ax.grid(axis="x", alpha=0.18)
    ax.text(summary["training_Na_Li_max"] + 5, 0, "5-25", va="center", fontsize=10, color="#16202D")
    ax.text(float(nominal["Na_Li_mass_ratio"]) - 8, 1, "381.55", va="center", ha="right", fontsize=10, color="#16202D")
    ax.set_title("Extrapolation boundary carried into every surrogate row", loc="left", fontsize=14, fontweight="bold")
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    return out


def render_nominal_transfer_card(nominal: dict[str, str]) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    out = ASSET_DIR / "nominal_transfer_math_card.png"
    fig, ax = plt.subplots(figsize=(8.6, 3.1), dpi=220)
    fig.patch.set_facecolor("#FFFFFF")
    ax.set_axis_off()
    ax.text(0.03, 0.86, "Nominal MS-2 transfer variables", fontsize=15, fontweight="bold", color="#16202D")
    ax.text(0.03, 0.66, r"$D_i=\left(E_i/(100-E_i)\right)/(O/A)$", fontsize=19, color="#16202D")
    rows = [
        (r"$E_{\mathrm{Li}}$", f"{float(nominal['li_extraction_pct']):.2f}%"),
        (r"$E_{\mathrm{Na}}$", f"{float(nominal['na_extraction_pct']):.2f}%"),
        (r"$S_{\mathrm{Li/Na}}$", f"{float(nominal['selectivity_Li_Na']):.2f}"),
        (r"$D_{\mathrm{Li}}$", f"{float(nominal['D_Li_phase_ratio_corrected']):.4f}"),
        (r"$D_{\mathrm{Na}}$", f"{float(nominal['D_Na_phase_ratio_corrected']):.4f}"),
    ]
    x0 = 0.04
    for i, (k, v) in enumerate(rows):
        x = x0 + i * 0.19
        ax.text(x, 0.37, k, fontsize=18, ha="center", color="#156576")
        ax.text(x, 0.18, v, fontsize=15, ha="center", fontweight="bold", color="#16202D")
    ax.text(0.03, 0.03, "O/A = 1.0, Li = 168 mg/L, Na = 64,100 mg/L", fontsize=9.5, color="#526172")
    fig.savefig(out, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    return out


def make_assets(summary: dict, nominal: dict[str, str]) -> dict[str, Path]:
    return {
        "pipeline": render_pipeline_diagram(summary),
        "process": render_process_diagram(),
        "validation": render_validation_chart(summary, nominal),
        "transfer": render_nominal_transfer_card(nominal),
        "formula_distribution": render_math_png("distribution_definition.png", r"$D_i=\left(E_i/(100-E_i)\right)/(O/A)$"),
        "formula_selectivity": render_math_png("selectivity_definition.png", r"$S_{\mathrm{Li/Na}}=D_{\mathrm{Li}}/D_{\mathrm{Na}}$"),
    }


def build_slides() -> list[dict]:
    summary = read_summary()
    nominal = read_nominal_row()
    assets = make_assets(summary, nominal)
    li_range = f"{summary['li_extraction_min']:.1f}-{summary['li_extraction_max']:.1f}%"
    na_range = f"{summary['na_extraction_min']:.1f}-{summary['na_extraction_max']:.1f}%"
    return [
        {
            "title": "Produced-Water Lithium Case Study",
            "subtitle": "Calibrated Rezaee DES/TOPO bridge for a pretreated Smackover Li/Na stream",
            "images": [
                asset("analyses/brine_composition_screening/results/figures/brine_li_tds_scatter.png"),
                asset("analyses/brine_composition_screening/results/figures/produced_water_candidate_scores.png"),
            ],
            "bullets": [
                "Use Smackover as the source brine because it combines high lithium grade with hypersaline operating severity.",
                "Start extraction after divalent pretreatment; the solvent model sees a clean Li/Na stream.",
                "Use Rezaee DES/TOPO as the active Li/Na chemistry bridge for the current presentation chain.",
            ],
        },
        {
            "title": "Pretreated Feed And Design Space",
            "subtitle": "The separation problem is posed after removing divalent cation interference",
            "images": [
                asset("analyses/rezaee_2026_pcsaft_epcsaft/results/surrogate_input_space/figures/rezaee_surrogate_input_ranges.png"),
                assets["process"],
            ],
            "bullets": [
                "Nominal clean feed: 168 mg/L Li and 64,100 mg/L Na from the MS-2 Smackover row.",
                "TDS is retained as a process severity feature, while Ca/Mg/Sr/Ba remain pretreatment variables.",
                "The UQ matrix varies TDS, lithium grade, and organic/aqueous ratio at fixed Rezaee chemistry settings.",
            ],
        },
        {
            "title": "Surrogate Construction",
            "subtitle": "A calibrated Section 3.2 DOE-basis surface supplies process-facing distribution coefficients",
            "images": [
                assets["pipeline"],
                assets["formula_distribution"],
            ],
            "bullets": [
                "Fit lightweight response surfaces in log-distribution space from 31 upstream Rezaee validation rows.",
                "Evaluate 523 produced-water rows with sodium derived from the selected Smackover TDS ratio.",
                "Carry the model basis and validity flag with every row used downstream.",
            ],
        },
        {
            "title": "Lithium Extraction Surface",
            "subtitle": "The calibrated surrogate returns finite extraction responses across the full screening matrix",
            "images": [
                asset("analyses/rezaee_2026_pcsaft_epcsaft/results/uq_surrogate/figures/calibrated_li_extraction_surface.png"),
            ],
            "bullets": [
                f"Lithium extraction spans {li_range}; sodium extraction spans {na_range}.",
                f"Nominal MS-2 one-stage lithium extraction is {float(nominal['li_extraction_pct']):.2f}%.",
                "The organic/aqueous ratio is the dominant process lever in the current three-variable sweep.",
            ],
        },
        {
            "title": "Selectivity And Transfer Variables",
            "subtitle": "The output is a staged-contacting contract, not only a recovery percentage",
            "images": [
                asset("analyses/rezaee_2026_pcsaft_epcsaft/results/uq_surrogate/figures/calibrated_selectivity_vs_na_li.png"),
                assets["transfer"],
            ],
            "bullets": [
                f"Nominal Li/Na selectivity is {float(nominal['selectivity_Li_Na']):.2f} at the clean MS-2 point.",
                "The corrected distribution coefficients are the primary variables for MSContactor handoff.",
                "Selectivity remains favorable under produced-water sodium loading, but it is extrapolated beyond the original DOE range.",
            ],
        },
        {
            "title": "PrOMMiS/IDAES Handoff",
            "subtitle": "Transfer variables are connected to staged recovery and screening economics",
            "images": [
                assets["process"],
                asset("analyses/rezaee_2026_pcsaft_epcsaft/results/uq_surrogate/figures/calibrated_costing_scenarios.png"),
            ],
            "bullets": [
                "The handoff table carries extraction, selectivity, distribution coefficients, raffinate, and extract loading.",
                "Costing scenarios use calibrated surrogate recovery values at the same feed-flow basis.",
                "Economics are screening-level until residence time, solvent loss, reagent demand, and product pricing are specified.",
            ],
        },
        {
            "title": "Validation Boundary",
            "subtitle": "The strongest caveat is the sodium-load extrapolation from Rezaee DOE space to Smackover water",
            "images": [
                assets["validation"],
                assets["formula_selectivity"],
            ],
            "bullets": [
                "Every produced-water row is flagged as outside the Rezaee DOE Na/Li training domain.",
                "The flag limits interpretation: process screening is appropriate; definitive high-Na/Li validation remains future work.",
                "The direct reactive-LLE closure is not claimed; the calibrated distribution surrogate is the presentation basis.",
            ],
        },
        {
            "title": "Presentation-Ready Chain",
            "subtitle": "The current story is now Rezaee-calibrated from feed definition through process screening",
            "images": [
                asset("analyses/rezaee_2026_pcsaft_epcsaft/results/uq_surrogate/figures/calibrated_li_extraction_surface.png"),
                asset("analyses/rezaee_2026_pcsaft_epcsaft/results/uq_surrogate/figures/calibrated_costing_scenarios.png"),
            ],
            "bullets": [
                "Use the calibrated Rezaee surrogate matrix as the active data product for the deck.",
                "Keep archived solvent-selection material outside this active Rezaee-calibrated presentation chain.",
                "Next technical step: replace algebraic process handoff with a full PrOMMiS MSContactor solve and costing block.",
            ],
        },
    ]


def add_text_box(slide, x, y, w, h, text_value, size=18, bold=False, color=BODY_COLOR, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text_value
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    return shape


def add_image_fit(slide, img_path: Path, x, y, w, h):
    with Image.open(img_path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        width = w
        height = w / img_ratio
        xx = x
        yy = y + (h - height) / 2
    else:
        height = h
        width = h * img_ratio
        xx = x + (w - width) / 2
        yy = y
    slide.shapes.add_picture(str(img_path), Inches(xx), Inches(yy), width=Inches(width), height=Inches(height))


def add_figure_surface(slide, img_path: Path, x, y, w, h):
    add_rect(slide, x, y, w, h, WHITE, None, False)
    add_image_fit(slide, img_path, x + 0.16, y + 0.15, w - 0.32, h - 0.30)


def add_takeaway_panel(slide, bullets: list[str], x=8.05, y=1.30, w=4.42, h=5.26):
    add_rect(slide, x, y, 0.06, h, ACCENT, None, False)
    add_text_box(slide, x + 0.28, y + 0.02, w - 0.32, 0.26, "ENGINEERING READOUT", 8.5, True, ACCENT)
    tx = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.42), Inches(w - 0.32), Inches(h - 0.50))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1:02d}  {bullet}"
        p.level = 0
        p.font.size = Pt(16.5)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(15)


def add_header(slide, idx: int, title: str, subtitle: str):
    add_rect(slide, 0, 0, 13.333, 0.10, ACCENT, None, False)
    add_text_box(slide, 0.55, 0.30, 8.8, 0.44, title, 24, True, DEEP_NAVY)
    add_text_box(slide, 0.57, 0.78, 8.2, 0.32, subtitle, 11.8, False, MUTED_COLOR)


def add_footer(slide):
    add_rect(slide, 0.55, 7.10, 1.35, 0.035, ACCENT, None, False)
    add_text_box(slide, 2.05, 6.99, 8.6, 0.25, "Calibrated Rezaee DES/TOPO clean Li/Na produced-water case study", 8.3, False, MUTED_COLOR)


def add_metric_strip(slide, metrics: list[tuple[str, str]], x, y, w):
    if not metrics:
        return
    gap = 0.12
    cell_w = (w - gap * (len(metrics) - 1)) / len(metrics)
    for i, (label, value) in enumerate(metrics):
        xx = x + i * (cell_w + gap)
        add_rect(slide, xx, y, cell_w, 0.66, TEAL_LIGHT if i == 0 else WHITE, None, False)
        add_text_box(slide, xx + 0.12, y + 0.10, cell_w - 0.24, 0.17, label.upper(), 7.5, True, ACCENT)
        add_text_box(slide, xx + 0.12, y + 0.31, cell_w - 0.24, 0.24, value, 13.5, True, DEEP_NAVY)


def delete_all_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        slide_id_list.remove(slide_id)


def build_pptx(slides: list[dict], out_pptx: Path) -> None:
    prs = Presentation(str(TEMPLATE_PATH)) if TEMPLATE_PATH.exists() else Presentation()
    if TEMPLATE_PATH.exists():
        delete_all_slides(prs)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0]
    for idx, spec in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PAPER_BG
        images = spec["images"]
        if idx == 1:
            add_rect(slide, 0, 0, 13.333, 7.5, WHITE, None, False)
            add_rect(slide, 0, 0, 0.30, 7.5, ACCENT, None, False)
            add_rect(slide, 0.30, 0, 0.08, 7.5, AMBER, None, False)
            add_text_box(slide, 0.72, 0.58, 5.8, 0.95, spec["title"], 31, True, DEEP_NAVY)
            add_text_box(slide, 0.76, 1.56, 5.35, 0.60, spec["subtitle"], 15.5, False, MUTED_COLOR)
            add_metric_strip(slide, [("Li feed", "168 mg/L"), ("Na feed", "64,100 mg/L"), ("Nominal Li extraction", "50.66%")], 0.76, 6.04, 5.75)
            add_figure_surface(slide, images[0], 6.85, 0.72, 5.70, 2.85)
            add_figure_surface(slide, images[1], 6.85, 3.88, 5.70, 2.50)
            add_takeaway_panel(slide, spec["bullets"], x=0.76, y=2.48, w=5.75, h=2.95)
            add_text_box(slide, 12.03, 0.34, 0.68, 0.25, f"{idx:02d}", 10, True, ACCENT, PP_ALIGN.RIGHT)
        else:
            add_header(slide, idx, spec["title"], spec["subtitle"])
            add_rect(slide, 0.55, 1.34, 6.98, 5.35, WHITE, None, False)
            if len(images) == 1:
                add_image_fit(slide, images[0], 0.78, 1.58, 6.52, 4.86)
            else:
                add_image_fit(slide, images[0], 0.80, 1.50, 6.45, 2.30)
                add_image_fit(slide, images[1], 0.80, 4.15, 6.45, 2.05)
            add_takeaway_panel(slide, spec["bullets"])
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)


def draw_image_fit(c: canvas.Canvas, img_path: Path, x, y, w, h):
    with Image.open(img_path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        width = w
        height = w / img_ratio
        xx = x
        yy = y + (h - height) / 2
    else:
        height = h
        width = h * img_ratio
        xx = x + (w - width) / 2
        yy = y
    c.drawImage(str(img_path), xx, yy, width=width, height=height, preserveAspectRatio=True, mask="auto")


def wrap_text(text_value: str, max_chars: int) -> list[str]:
    words = text_value.split()
    lines = []
    cur = ""
    for word in words:
        test = (cur + " " + word).strip()
        if len(test) > max_chars:
            lines.append(cur)
            cur = word
        else:
            cur = test
    if cur:
        lines.append(cur)
    return lines


def build_pdf(slides: list[dict], out_pdf: Path) -> None:
    page = landscape((13.333 * inch, 7.5 * inch))
    c = canvas.Canvas(str(out_pdf), pagesize=page)
    pw, ph = page
    for idx, spec in enumerate(slides, 1):
        c.setFillColor(colors.HexColor("#F8FAFC"))
        c.rect(0, 0, pw, ph, fill=1, stroke=0)
        c.setFillColor(colors.HexColor("#16202D"))
        c.setFont("Helvetica-Bold", 26)
        c.drawString(0.42 * inch, 6.97 * inch, spec["title"])
        c.setFillColor(colors.HexColor("#5D6C7E"))
        c.setFont("Helvetica", 12.5)
        c.drawString(0.45 * inch, 6.60 * inch, spec["subtitle"])
        c.setFont("Helvetica-Bold", 10)
        c.drawRightString(12.78 * inch, 6.98 * inch, f"{idx:02d}")
        imgs = spec["images"]
        if len(imgs) == 1:
            draw_image_fit(c, imgs[0], 0.45 * inch, 0.78 * inch, 7.25 * inch, 5.35 * inch)
        else:
            draw_image_fit(c, imgs[0], 0.45 * inch, 3.62 * inch, 7.25 * inch, 2.55 * inch)
            draw_image_fit(c, imgs[1], 0.45 * inch, 0.78 * inch, 7.25 * inch, 2.62 * inch)
        c.setFillColor(colors.white)
        c.setStrokeColor(colors.HexColor("#D1DAE6"))
        c.roundRect(8.10 * inch, 0.70 * inch, 4.72 * inch, 5.72 * inch, 6, fill=1, stroke=1)
        y = 5.95 * inch
        c.setFont("Helvetica", 18.5)
        c.setFillColor(colors.HexColor("#263242"))
        for bullet in spec["bullets"]:
            lines = wrap_text(bullet, 38)
            c.circle(8.34 * inch, y + 4, 2.1, fill=1, stroke=0)
            for line in lines:
                c.drawString(8.47 * inch, y, line)
                y -= 0.245 * inch
            y -= 0.19 * inch
        c.setFont("Helvetica", 8.5)
        c.setFillColor(colors.HexColor("#5D6C7E"))
        c.drawString(0.48 * inch, 0.18 * inch, "Calibrated Rezaee DES/TOPO clean Li/Na case study - produced-water process screening")
        c.showPage()
    c.save()


def build_quick_report(slides: list[dict], out_png: Path) -> None:
    font = ImageFont.load_default()
    cols = 2
    thumb = (560, 315)
    pad = 24
    label_h = 42
    title_h = 54
    rows = (len(slides) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * thumb[0] + (cols + 1) * pad, title_h + rows * (thumb[1] + label_h + pad) + pad), "white")
    draw = ImageDraw.Draw(sheet)
    draw.text((pad, 18), "Quick visual report - calibrated Rezaee case-study deck", fill=(20, 28, 36), font=font)
    for idx, spec in enumerate(slides):
        x = pad + (idx % cols) * (thumb[0] + pad)
        y = title_h + (idx // cols) * (thumb[1] + label_h + pad)
        frame = Image.new("RGB", thumb, (248, 250, 252))
        imgs = spec["images"]
        if len(imgs) == 1:
            im = Image.open(imgs[0]).convert("RGB")
            im.thumbnail((332, 248), Image.LANCZOS)
            frame.paste(im, (22, 48))
        else:
            for j, img in enumerate(imgs[:2]):
                im = Image.open(img).convert("RGB")
                im.thumbnail((288, 118), Image.LANCZOS)
                frame.paste(im, (20, 46 + j * 134))
        d = ImageDraw.Draw(frame)
        d.text((18, 12), spec["title"], fill=(20, 28, 36), font=font)
        by = 50
        for bullet in spec["bullets"][:3]:
            lines = wrap_text(bullet, 29)
            d.text((370, by), "- " + lines[0], fill=(30, 41, 59), font=font)
            by += 15
            for extra in lines[1:3]:
                d.text((380, by), extra, fill=(30, 41, 59), font=font)
                by += 15
            by += 9
        sheet.paste(frame, (x, y))
        draw.rectangle([x, y, x + thumb[0] - 1, y + thumb[1] - 1], outline=(190, 198, 208))
        draw.text((x, y + thumb[1] + 8), f"{idx + 1}. {spec['subtitle']}", fill=(84, 96, 112), font=font)
    sheet.save(out_png)


def main() -> None:
    slides = build_slides()
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUT_DIR / "Rezaee_Calibrated_Produced_Water_Case_Study_2026_05_08.pptx"
    pdf_path = OUT_DIR / "Rezaee_Calibrated_Produced_Water_Case_Study_2026_05_08.pdf"
    quick_png = OUT_DIR / "quick_visual_report.png"
    build_pptx(slides, pptx_path)
    if pdf_path.exists():
        pdf_path.unlink()
    build_quick_report(slides, quick_png)
    print(pptx_path)
    print("Export PDF/pages with scripts/presentation/export_rezaee_deck_with_powerpoint.ps1")
    print(quick_png)


if __name__ == "__main__":
    main()
